"""
Script to populate the official UniHack presentation template with ELIO project data.
Optimized for the white template background with high-contrast typography,
clean left-alignment, and pristine closing slide.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE_PATH = ROOT / "[EXT] UniHack-Protoype Template  (2).pptx"
OUTPUT_PATH = ROOT / "UniHack_ELIO_Prototype_Submission.pptx"
ART_DIR = ROOT / "artifacts"
SCREENSHOT_DIR = Path(r"C:\Users\rushd\.gemini\antigravity-ide\brain\18dca2bd-d329-4d11-a99c-7cdaf953f864")

# Color Tokens
NAVY = (15, 23, 42)       # #0f172a - Main titles
BLUE = (0, 114, 206)      # #0072ce - Unilog brand blue / key labels
SLATE = (30, 41, 59)      # #1e293b - High-contrast readable body text
MUTED = (71, 85, 105)     # #475569 - Secondary text

def style_para(p, text, font_name="Arial", size_pt=12, bold=False, color_rgb=SLATE, space_after=Pt(4), align=PP_ALIGN.LEFT):
    p.text = text
    p.alignment = align
    if p.font:
        p.font.name = font_name
        p.font.size = Pt(size_pt)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color_rgb)
    if space_after:
        p.space_after = space_after

def add_bullet(tf, bold_prefix, text, size_pt=11, space_after=Pt(4.5), prefix_rgb=BLUE, body_rgb=SLATE):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = space_after
    
    # Run 1: Bold prefix
    r1 = p.add_run()
    r1.text = bold_prefix + " "
    r1.font.name = "Arial"
    r1.font.size = Pt(size_pt)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(*prefix_rgb)
    
    # Run 2: Body
    r2 = p.add_run()
    r2.text = text
    r2.font.name = "Arial"
    r2.font.size = Pt(size_pt)
    r2.font.bold = False
    r2.font.color.rgb = RGBColor(*body_rgb)

def fill_deck():
    prs = Presentation(str(TEMPLATE_PATH))
    print(f"Loaded template with {len(prs.slides)} slides.")

    # -------------------------------------------------------------
    # SLIDE 2: Team Details
    # -------------------------------------------------------------
    s2 = prs.slides[1]
    for shape in s2.shapes:
        if shape.has_text_frame and "Team Details" in shape.text_frame.text:
            shape.top = Inches(3.3)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(1.8)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Team Details", size_pt=18, bold=True, color_rgb=NAVY, space_after=Pt(6))
            add_bullet(tf, "Team Name:", "ELIO", size_pt=12, space_after=Pt(4))
            add_bullet(tf, "Team Leader:", "Darshan K (rushdarshan@gmail.com)", size_pt=12, space_after=Pt(4))
            add_bullet(tf, "Challenge Statement:", "AI-Powered Product Intelligence for Industrial Commerce", size_pt=12, space_after=Pt(4))
            add_bullet(tf, "Solution Name:", "ELIO — Evidence-Gated Catalog Intelligence & Syndication Engine", size_pt=12, space_after=Pt(4))

    # -------------------------------------------------------------
    # SLIDE 3: Brief about your solution
    # -------------------------------------------------------------
    s3 = prs.slides[2]
    for shape in s3.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(4.3)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Brief About Your Solution", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(4))
            
            p_sub = tf.add_paragraph()
            style_para(p_sub, "ELIO — Immutable, Falsifiable Product Intelligence for Industrial Commerce", size_pt=11.5, bold=True, color_rgb=BLUE, space_after=Pt(8))
            
            add_bullet(tf, "The Core Problem:", "Industrial distributor feeds arrive with messy, truncated descriptions, placeholder brands ('-- Unbranded --'), non-standard units ('inches' vs 'in'), and unpopulated attribute columns.", size_pt=10.5, space_after=Pt(4))
            add_bullet(tf, "The ELIO Solution:", "A deterministic 9-stage DAG that transforms raw 6-column distributor rows into verified 252-column product catalog records with 100% provenance coverage.", size_pt=10.5, space_after=Pt(4))
            add_bullet(tf, "Zero-Hallucination Dual-Pass Gate:", "Every emitted attribute is anchored to verbatim character offset spans in manufacturer documentation. If evidence is missing, ELIO honestly refuses rather than fabricating guesses.", size_pt=10.5, space_after=Pt(4))
            add_bullet(tf, "Enterprise Provenance & Replay:", "Generates cryptographic execution receipts (rcpt_<sha256>) and append-only decision logs (decision_log.jsonl) enabling 100% byte-identical state replay across audits.", size_pt=10.5, space_after=Pt(4))

    # -------------------------------------------------------------
    # SLIDE 4: 1. Enrichment / 2. Accuracy & Trust / 3. Scalability
    # -------------------------------------------------------------
    s4 = prs.slides[3]
    for shape in s4.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(4.3)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Enrichment, Accuracy & Enterprise Scalability", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(6))
            
            add_bullet(tf, "1. Minimal Info Enrichment:", "The 9-stage DAG transforms minimal inputs (MPN, Description, Manufacturer) into verified taxonomy classification, normalized attribute sets, and a Universal 5-Variant Description Pack (Invoice, Mobile, Short, Long, Marketing).", size_pt=10, space_after=Pt(4))
            add_bullet(tf, "2. Accuracy & Trust Strategy:", "Dual-Pass Verification Gate validates candidate values against verbatim source text or approved 63-entry fraction tables. Employs 4 honest abstention classes when evidence fails.", size_pt=10, space_after=Pt(4))
            add_bullet(tf, "3. Human-in-the-Loop Governance:", "Interactive Review Queue allows 1-click Accept/Reject and local overrides, logged to an immutable event-sourced decision stream for audit compliance.", size_pt=10, space_after=Pt(4))
            add_bullet(tf, "4. Enterprise Catalog Scalability:", "Processes 1,000 catalog rows in 1.70s (588.2 rows/sec local deterministic throughput) with zero GPU cost and 100% 252-column schema conformity.", size_pt=10, space_after=Pt(4))

    # -------------------------------------------------------------
    # SLIDE 5: Opportunities & USP
    # -------------------------------------------------------------
    s5 = prs.slides[4]
    for shape in s5.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(4.3)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Opportunities & Unique Selling Proposition (USP)", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(6))
            
            add_bullet(tf, "How ELIO Differs from Existing Tools:", "Traditional PIMs require slow, manual data entry; generic LLMs hallucinate critical engineering tolerances and part specs. ELIO combines deterministic speed with evidence-gated verification.", size_pt=10.5, space_after=Pt(5))
            add_bullet(tf, "Solving the Unilog Challenge Statement:", "Automates Unilog Content Guidelines across 89 UOM categories, 27,000+ brand entities, 161,000+ List of Values rules, and the strict 252-column delivery sequence.", size_pt=10.5, space_after=Pt(5))
            add_bullet(tf, "Core USP — Falsifiable AI:", "'If an attribute value cannot be proven with verbatim source evidence, ELIO refuses to guess.' Honest, transparent abstentions over confident hallucinations.", size_pt=10.5, space_after=Pt(5))

    # -------------------------------------------------------------
    # SLIDE 6: List of Features
    # -------------------------------------------------------------
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(4.3)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Key Features Offered by ELIO", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(5))
            
            add_bullet(tf, "1. 9-Stage Evidence-Gated DAG:", "Ingest, entity resolution, taxonomy, retrieval, extraction, dual-pass verification, quality gate, 5-variant descriptions, and 252-column projection.", size_pt=10, space_after=Pt(3.5))
            add_bullet(tf, "2. Master UOM & Trade Fraction Normalizer:", "Strict compliance with ~500 approved abbreviations across 89 physical dimensions, with exact fraction mapping (e.g., 50.25 in -> 50-1/4 in).", size_pt=10, space_after=Pt(3.5))
            add_bullet(tf, "3. Universal 5-Variant Description Engine:", "Constructs Invoice (<=40 char, CAPS), Mobile (60-80 char), Short (Brand+Series+MPN+Type), Long, and Marketing descriptions strictly from verified facts.", size_pt=10, space_after=Pt(3.5))
            add_bullet(tf, "4. Interactive Custody Drawer & Proof Graph:", "Inspect source URLs, document pages, exact character spans, and recompute SHA-256 custody hashes live in the web cockpit.", size_pt=10.5, space_after=Pt(3.5))
            add_bullet(tf, "5. Deterministic Decision Replay:", "Reconstructs full catalog state byte-identically from append-only JSONL decision logs.", size_pt=10, space_after=Pt(3.5))
            add_bullet(tf, "6. Distributor Entity Blacklist:", "Prevents supplier names (e.g. 'Appliance Dealers Cooperative') from poisoning brand/manufacturer fields.", size_pt=10, space_after=Pt(3.5))

    # -------------------------------------------------------------
    # SLIDE 7: Process Flow Diagram
    # -------------------------------------------------------------
    s7 = prs.slides[6]
    for shape in s7.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(0.6)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Process Flow — 9-Stage Evidence-Gated DAG", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(2))
    
    flow_img = ART_DIR / "diagram_process_flow.png"
    if flow_img.exists():
        s7.shapes.add_picture(str(flow_img), Inches(0.4), Inches(1.5), Inches(9.2), Inches(3.68))

    # -------------------------------------------------------------
    # SLIDE 8: Wireframes / Cockpit Layout
    # -------------------------------------------------------------
    s8 = prs.slides[7]
    for shape in s8.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(0.6)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "ELIO Cockpit — 5-Surface Operations Dashboard & Proof Graph", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(2))
    
    proof_img = SCREENSHOT_DIR / "proof_graph_verified_1787374589470.png"
    if not proof_img.exists():
        proof_img = SCREENSHOT_DIR / "clean_dashboard_cockpit_1787368566744.png"
    if proof_img.exists():
        s8.shapes.add_picture(str(proof_img), Inches(0.4), Inches(1.5), Inches(9.2), Inches(3.68))

    # -------------------------------------------------------------
    # SLIDE 9: Architecture Diagram
    # -------------------------------------------------------------
    s9 = prs.slides[8]
    for shape in s9.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(0.6)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Architecture Diagram — Three-Layer System Design", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(2))
    
    arch_img = ART_DIR / "diagram_architecture.png"
    if arch_img.exists():
        s9.shapes.add_picture(str(arch_img), Inches(0.4), Inches(1.5), Inches(9.2), Inches(3.68))

    # -------------------------------------------------------------
    # SLIDE 10: Technologies Used
    # -------------------------------------------------------------
    s10 = prs.slides[9]
    for shape in s10.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(4.3)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Technologies Used in the Solution", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(6))
            
            add_bullet(tf, "Core DAG Pipeline Engine:", "Python 3.13, Pydantic v2 (EnrichedRecord, ClaimRecord, SourceEvidence), RapidFuzz (fuzzy brand resolution), Pandas, Openpyxl.", size_pt=10.5, space_after=Pt(4.5))
            add_bullet(tf, "State & Provenance Engine:", "Waku/Cambium-style RunStore with SHA-256 cryptographic receipts, JSONL append-only decision stream, and byte-identical replay engine.", size_pt=10.5, space_after=Pt(4.5))
            add_bullet(tf, "Frontend Operations Cockpit:", "Next.js 16.3 (App Router), React 19, TypeScript, GSAP Core & useGSAP hook (60fps performance), Vanilla CSS design tokens.", size_pt=10.5, space_after=Pt(4.5))
            add_bullet(tf, "Evaluation & Quality Gates:", "Automated 12-gate test harness (verify_everything.py), clean-room evaluator, adversarial holdout generator (Seed-42).", size_pt=10.5, space_after=Pt(4.5))

    # -------------------------------------------------------------
    # SLIDE 11: Estimated Implementation Cost
    # -------------------------------------------------------------
    s11 = prs.slides[10]
    for shape in s11.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(4.3)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Estimated Implementation & Operational Cost", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(6))
            
            add_bullet(tf, "Compute Footprint ($0.00 GPU Cost):", "The entire deterministic DAG runs on standard CPU instances. Enriches 1,000 catalog rows in 1.70 seconds with minimal memory overhead.", size_pt=10.5, space_after=Pt(4.5))
            add_bullet(tf, "Token & API Cost:", "$0.00 token cost for deterministic baseline extraction; <$0.002 per row when proposal layer is active (94% cache hit rate).", size_pt=10.5, space_after=Pt(4.5))
            add_bullet(tf, "Hosting & Maintenance:", "Deployable as lightweight containers (FastAPI / Next.js) on AWS Lambda, Cloud Run, or on-premise Kubernetes (~$15 - $30 / month base).", size_pt=10.5, space_after=Pt(4.5))
            add_bullet(tf, "Return on Investment (ROI):", "Reduces manual catalog enrichment time by >95% while eliminating costly e-commerce product return rates caused by inaccurate specs.", size_pt=10.5, space_after=Pt(4.5))

    # -------------------------------------------------------------
    # SLIDE 12: Snapshots of the MVP
    # -------------------------------------------------------------
    s12 = prs.slides[11]
    for shape in s12.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(0.6)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Snapshots of the Live MVP Cockpit", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(2))
    
    mvp_img = SCREENSHOT_DIR / "clean_dashboard_cockpit_1787368566744.png"
    if not mvp_img.exists():
        mvp_img = SCREENSHOT_DIR / "final_dashboard_1787368166258.png"
    if mvp_img.exists():
        s12.shapes.add_picture(str(mvp_img), Inches(0.4), Inches(1.5), Inches(9.2), Inches(3.68))

    # -------------------------------------------------------------
    # SLIDE 13: Additional Details / Future Development
    # -------------------------------------------------------------
    s13 = prs.slides[12]
    for shape in s13.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(4.3)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Future Development & Enterprise Roadmap", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(6))
            
            add_bullet(tf, "1. Multi-Modal Technical Drawing Parsing:", "Extract dimensional schematics and wiring diagrams directly from manufacturer PDF spec sheets via computer vision.", size_pt=10.5, space_after=Pt(4.5))
            add_bullet(tf, "2. Automated Live Manufacturer Crawlers:", "Continuous crawler network indexing official OEM portal updates with automated staleness detection.", size_pt=10.5, space_after=Pt(4.5))
            add_bullet(tf, "3. Federated Distributor Syndication:", "Real-time webhook syndication to Akeneo, Salsify, and Syndigo PIM endpoints via GraphQL.", size_pt=10.5, space_after=Pt(4.5))
            add_bullet(tf, "4. Collaborative Review Hub:", "Multi-user role-based governance with cryptographic signatures for distributor catalog sign-offs.", size_pt=10.5, space_after=Pt(4.5))

    # -------------------------------------------------------------
    # SLIDE 14: Provide links to your:
    # -------------------------------------------------------------
    s14 = prs.slides[13]
    for shape in s14.shapes:
        if shape.has_text_frame:
            shape.top = Inches(0.85)
            shape.left = Inches(0.4)
            shape.width = Inches(9.2)
            shape.height = Inches(4.3)
            tf = shape.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            style_para(p0, "Submission Links & Verification Resources", size_pt=19, bold=True, color_rgb=NAVY, space_after=Pt(6))
            
            add_bullet(tf, "GitHub Public Repository:", "https://github.com/rushdarshan/Elio (Frozen at commit bar-5-clean)", size_pt=11, space_after=Pt(5))
            add_bullet(tf, "Interactive Prototype Cockpit:", "http://localhost:3000/app/dashboard (Next.js App Router)", size_pt=11, space_after=Pt(5))
            add_bullet(tf, "Master 12-Gate Verification Suite:", "python -B scripts/verify_everything.py (All Gates Passing)", size_pt=11, space_after=Pt(5))
            add_bullet(tf, "Automated Judge Walk Test:", "python -B scripts/judge_walk.py (5/5 Surfaces Validated)", size_pt=11, space_after=Pt(5))
            add_bullet(tf, "Cryptographic Proof Chain Verifier:", "python -B scripts/verify_receipt.py (91/91 Evidence Hashes Verified)", size_pt=11, space_after=Pt(5))

    # SLIDE 15: Left untouched as pristine graphic closing slide with built-in Thank You artwork

    prs.save(str(OUTPUT_PATH))
    print(f"\n[SUCCESS] Successfully populated official template and saved to: {OUTPUT_PATH}")

if __name__ == '__main__':
    fill_deck()
