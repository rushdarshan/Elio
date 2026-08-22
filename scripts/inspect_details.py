from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation('[EXT] UniHack-Protoype Template  (2).pptx')
for idx, slide in enumerate(prs.slides):
    print(f"\n==================== SLIDE {idx+1} ====================")
    for s_idx, shape in enumerate(slide.shapes):
        print(f"Shape {s_idx}: name='{shape.name}', type={shape.shape_type}, left={shape.left.inches:.2f}\", top={shape.top.inches:.2f}\", width={shape.width.inches:.2f}\", height={shape.height.inches:.2f}\"")
        if shape.has_text_frame:
            tf = shape.text_frame
            for p_idx, p in enumerate(tf.paragraphs):
                font_name = p.font.name if p.font else None
                font_size = p.font.size.pt if p.font and p.font.size else None
                bold = p.font.bold if p.font else None
                print(f"  P{p_idx} (font={font_name}, size={font_size}, bold={bold}): {p.text}")
