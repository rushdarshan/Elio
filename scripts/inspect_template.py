from pptx import Presentation

prs = Presentation('[EXT] UniHack-Protoype Template  (2).pptx')
print(f'Total slides: {len(prs.slides)}')
for idx, slide in enumerate(prs.slides):
    print(f'\n==================== SLIDE {idx+1} ====================')
    for s_idx, shape in enumerate(slide.shapes):
        has_text = shape.has_text_frame
        text_preview = ""
        if has_text:
            text_preview = " | ".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
        print(f'  Shape {s_idx} [{shape.name}] (type: {shape.shape_type}): {text_preview[:120]}')
